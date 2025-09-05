"""
Document processing service for extracting text from various file formats.
"""
import re
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional
from abc import ABC, abstractmethod

# Third-party imports
try:
    import PyPDF2
    import pandas as pd
    from docx import Document
    from pptx import Presentation
    import markdown
except ImportError as e:
    raise ImportError(f"Missing required package: {e}")


class DocumentProcessor:
    """Enterprise-grade document processing engine."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_extensions = {'.pdf', '.txt', '.doc', '.docx', '.csv', '.xlsx', '.xls', '.md', '.pptx'}
        
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        Process a file and extract its text content.
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            Dictionary containing file metadata and extracted text
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        # Extract text based on file type
        text = self._extract_text_by_type(file_path, file_extension)
        
        # Create intelligent chunks
        chunks = self._create_chunks(text)
        
        return {
            'file_path': str(file_path),
            'file_name': file_path.name,
            'file_type': file_extension,
            'text': text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'file_size': file_path.stat().st_size
        }
    
    def _extract_text_by_type(self, file_path: Path, file_extension: str) -> str:
        """Extract text based on file type."""
        extractors = {
            '.pdf': self._extract_pdf_text,
            '.txt': self._extract_txt_text,
            '.doc': self._extract_docx_text,
            '.docx': self._extract_docx_text,
            '.csv': self._extract_excel_text,
            '.xlsx': self._extract_excel_text,
            '.xls': self._extract_excel_text,
            '.md': self._extract_markdown_text,
            '.pptx': self._extract_pptx_text
        }
        
        extractor = extractors.get(file_extension)
        if not extractor:
            raise ValueError(f"No extractor available for {file_extension}")
            
        return extractor(file_path)
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error reading PDF {file_path}: {str(e)}")
        return text.strip()
    
    def _extract_txt_text(self, file_path: Path) -> str:
        """Extract text from TXT file."""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        raise Exception(f"Could not decode text file {file_path}")
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading DOCX {file_path}: {str(e)}")
    
    def _extract_excel_text(self, file_path: Path) -> str:
        """Extract text from Excel/CSV file."""
        try:
            if file_path.suffix.lower() == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            
            text = f"Document: {file_path.name}\n\n"
            text += df.to_string(index=False)
            return text
        except Exception as e:
            raise Exception(f"Error reading Excel/CSV {file_path}: {str(e)}")
    
    def _extract_markdown_text(self, file_path: Path) -> str:
        """Extract text from Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert markdown to plain text
            html = markdown.markdown(md_content)
            text = re.sub('<[^<]+?>', '', html)
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading Markdown {file_path}: {str(e)}")
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        try:
            prs = Presentation(file_path)
            text = f"Presentation: {file_path.name}\n\n"
            
            for i, slide in enumerate(prs.slides, 1):
                text += f"Slide {i}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading PPTX {file_path}: {str(e)}")
    
    def _create_chunks(self, text: str) -> List[Dict[str, Any]]:
        """Create intelligent text chunks with overlap."""
        if not text.strip():
            return []
        
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        chunks = []
        current_chunk = ""
        current_size = 0
        
        for sentence in sentences:
            sentence_size = len(sentence)
            
            if current_size + sentence_size > self.chunk_size and current_chunk:
                chunks.append({
                    'text': current_chunk.strip(),
                    'size': current_size,
                    'chunk_id': len(chunks)
                })
                
                overlap_text = self._get_overlap_text(current_chunk)
                current_chunk = overlap_text + " " + sentence
                current_size = len(current_chunk)
            else:
                current_chunk += " " + sentence if current_chunk else sentence
                current_size += sentence_size + (1 if current_chunk else 0)
        
        if current_chunk.strip():
            chunks.append({
                'text': current_chunk.strip(),
                'size': current_size,
                'chunk_id': len(chunks)
            })
        
        return chunks
    
    def _get_overlap_text(self, text: str) -> str:
        """Get intelligent overlap text from chunk end."""
        if len(text) <= self.chunk_overlap:
            return text
        
        overlap_start = len(text) - self.chunk_overlap
        overlap_text = text[overlap_start:]
        
        # Find sentence boundary
        sentence_break = overlap_text.find('. ')
        if sentence_break != -1:
            return overlap_text[sentence_break + 2:]
        
        # Find word boundary
        word_break = overlap_text.find(' ')
        if word_break != -1:
            return overlap_text[word_break + 1:]
        
        return overlap_text